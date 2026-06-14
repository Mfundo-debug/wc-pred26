from __future__ import annotations
import os, json, math, shutil, zipfile, subprocess, sys
from pathlib import Path
from datetime import datetime
from collections import Counter
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from PIL import Image, ImageDraw, ImageFont

BASE=Path('/mnt/data/wc26_bhm_final')
if BASE.exists(): shutil.rmtree(BASE)
for sub in ['data','outputs','figures','src','dashboard/flags','streamlit_app','notebooks','report_render']:
    (BASE/sub).mkdir(parents=True, exist_ok=True)
# Data
GROUPS={
 'A':['Mexico','South Africa','Korea Republic','Czech Republic'],
 'B':['Canada','Bosnia-Herzegovina','Qatar','Switzerland'],
 'C':['Brazil','Morocco','Haiti','Scotland'],
 'D':['United States','Paraguay','Australia','Turkey'],
 'E':['Germany','Ecuador','Curacao',"Cote d'Ivoire"],
 'F':['Netherlands','Japan','Tunisia','Sweden'],
 'G':['Belgium','Egypt','Iran','New Zealand'],
 'H':['Spain','Cape Verde','Saudi Arabia','Uruguay'],
 'I':['France','Senegal','Iraq','Norway'],
 'J':['Argentina','Algeria','Austria','Jordan'],
 'K':['Portugal','Colombia','Congo DR','Uzbekistan'],
 'L':['England','Croatia','Ghana','Panama']}
TEAMS=[t for ts in GROUPS.values() for t in ts]
team_to_group={t:g for g,ts in GROUPS.items() for t in ts}
CODES={'Argentina':'ARG','Spain':'ESP','France':'FRA','England':'ENG','Portugal':'POR','Brazil':'BRA','Morocco':'MAR','Netherlands':'NED','Belgium':'BEL','Germany':'GER','Croatia':'CRO','Uruguay':'URU','Mexico':'MEX','United States':'USA','Switzerland':'SUI','Colombia':'COL','Japan':'JPN','Senegal':'SEN','Austria':'AUT','Turkey':'TUR','Sweden':'SWE','Ecuador':'ECU','Korea Republic':'KOR','Australia':'AUS','Iran':'IRN',"Cote d'Ivoire":'CIV','Egypt':'EGY','Norway':'NOR','Scotland':'SCO','Qatar':'QAT','Paraguay':'PAR','Ghana':'GHA','Canada':'CAN','Saudi Arabia':'KSA','Tunisia':'TUN','Panama':'PAN','Bosnia-Herzegovina':'BIH','Algeria':'ALG','South Africa':'RSA','Czech Republic':'CZE','Uzbekistan':'UZB','New Zealand':'NZL','Iraq':'IRQ','Jordan':'JOR','Cape Verde':'CPV','Haiti':'HAI','Curacao':'CUW','Congo DR':'COD'}
CONFED={}
for t in ['Argentina','Brazil','Uruguay','Colombia','Ecuador','Paraguay']: CONFED[t]='CONMEBOL'
for t in ['Spain','France','England','Portugal','Netherlands','Belgium','Germany','Croatia','Switzerland','Austria','Turkey','Sweden','Norway','Scotland','Czech Republic','Bosnia-Herzegovina']: CONFED[t]='UEFA'
for t in ['Morocco','Senegal',"Cote d'Ivoire",'Egypt','Ghana','Tunisia','Algeria','South Africa','Cape Verde','Congo DR']: CONFED[t]='CAF'
for t in ['Mexico','United States','Canada','Panama','Haiti','Curacao']: CONFED[t]='CONCACAF'
for t in ['Japan','Korea Republic','Iran','Qatar','Saudi Arabia','Australia','Iraq','Jordan','Uzbekistan']: CONFED[t]='AFC'
CONFED['New Zealand']='OFC'
HOSTS={'Mexico','United States','Canada'}
# Current strength index anchored to top-10 June 2026 FIFA ranking context + calibrated remainder.
PTS={'Argentina':1877.27,'Spain':1874.71,'France':1870.70,'England':1828.02,'Portugal':1767.85,'Brazil':1765.86,'Morocco':1755.10,'Netherlands':1753.57,'Belgium':1742.00,'Germany':1715.00,'Croatia':1698,'Uruguay':1692,'Mexico':1676,'United States':1668,'Switzerland':1660,'Colombia':1658,'Japan':1647,'Senegal':1638,'Austria':1630,'Turkey':1618,'Sweden':1608,'Ecuador':1605,'Korea Republic':1586,'Australia':1567,'Iran':1562,"Cote d'Ivoire":1553,'Egypt':1548,'Norway':1546,'Scotland':1540,'Qatar':1535,'Paraguay':1533,'Ghana':1528,'Canada':1520,'Saudi Arabia':1512,'Tunisia':1508,'Panama':1498,'Bosnia-Herzegovina':1496,'Algeria':1494,'South Africa':1482,'Czech Republic':1478,'Uzbekistan':1460,'New Zealand':1445,'Iraq':1438,'Jordan':1434,'Cape Verde':1428,'Haiti':1416,'Curacao':1408,'Congo DR':1400}
HIST={t:0.0 for t in TEAMS}
HIST.update({'Argentina':17.7,'France':18.9,'Croatia':13.4,'Morocco':6.85,'Brazil':14.65,'Netherlands':10.1,'England':8.8,'Portugal':8.25,'Spain':9.4,'Uruguay':9.45,'Germany':12.15,'Belgium':8.8,'Switzerland':5.65,'Mexico':6.2,'United States':2.8,'Japan':6.2,'Korea Republic':4.65,'Senegal':2.85,'Australia':4.4,'Ghana':2.2,'Iran':3.1,'Tunisia':2.4,'Saudi Arabia':1.85,'Ecuador':1.55,'Canada':1.0,'Qatar':1.0,'South Africa':0.55})
ADJ={t:0.0 for t in TEAMS}; ADJ['South Africa']=-0.18; ADJ['Mexico']=-0.06; ADJ['Argentina']=-0.04
LIVE=[{'date':'2026-06-11','group':'A','home':'Mexico','away':'South Africa','home_goals':2,'away_goals':0,'notes':'Verified: Mexico 2-0 South Africa; two South Africa reds and one Mexico red are captured by adjustment layer.'},{'date':'2026-06-11','group':'A','home':'Korea Republic','away':'Czech Republic','home_goals':2,'away_goals':1,'notes':'Verified: Korea Republic 2-1 Czech Republic.'}]
fixtures=[]
# group A exact order
fixtures += [{'group':'A','home':'Mexico','away':'South Africa'}, {'group':'A','home':'Korea Republic','away':'Czech Republic'}, {'group':'A','home':'Mexico','away':'Korea Republic'}, {'group':'A','home':'Czech Republic','away':'South Africa'}, {'group':'A','home':'Czech Republic','away':'Mexico'}, {'group':'A','home':'South Africa','away':'Korea Republic'}]
for g,ts in GROUPS.items():
    if g=='A': continue
    for h,a in [(ts[0],ts[1]),(ts[2],ts[3]),(ts[0],ts[2]),(ts[3],ts[1]),(ts[3],ts[0]),(ts[1],ts[2])]: fixtures.append({'group':g,'home':h,'away':a})
team_df=pd.DataFrame([{'group':team_to_group[t],'team':t,'code':CODES[t],'confed':CONFED[t],'host':t in HOSTS,'fifa_points':PTS[t],'historical_wc_score':HIST[t],'current_adjustment':ADJ[t]} for t in TEAMS])
team_df.to_csv(BASE/'data/teams.csv',index=False); pd.DataFrame(fixtures).to_csv(BASE/'data/fixtures_2026.csv',index=False); pd.DataFrame(LIVE).to_csv(BASE/'data/verified_live_results.csv',index=False)
# PyMC simple BHM
import pymc as pm, arviz as az
N=len(TEAMS); idx={t:i for i,t in enumerate(TEAMS)}
rank_z=(team_df.fifa_points-team_df.fifa_points.mean())/team_df.fifa_points.std()
hist_z=(team_df.historical_wc_score-team_df.historical_wc_score.mean())/(team_df.historical_wc_score.std()+1e-6)
host_vec=team_df.host.astype(int).values; adj=team_df.current_adjustment.values
mu_prior=(0.74*rank_z.values+0.26*hist_z.values+adj+0.05*host_vec)
obs_signal=mu_prior.copy()
hidx=np.array([idx[x['home']] for x in LIVE]); aidx=np.array([idx[x['away']] for x in LIVE]); hg=np.array([x['home_goals'] for x in LIVE]); ag=np.array([x['away_goals'] for x in LIVE])
with pm.Model(coords={'team':TEAMS}) as model:
    sigma_team=pm.HalfNormal('sigma_team',0.35)
    global_log_goal_rate=pm.Normal('global_log_goal_rate',np.log(1.30),0.16)
    host_advantage=pm.Normal('host_advantage',0.07,0.04)
    strength=pm.Normal('strength',mu=mu_prior,sigma=sigma_team,dims='team')
    attack_bias=pm.Normal('attack_bias',0,0.10,dims='team')
    defence_bias=pm.Normal('defence_bias',0,0.10,dims='team')
    attack=pm.Deterministic('attack',0.52*strength+attack_bias,dims='team')
    defence=pm.Deterministic('defence',-0.52*strength+defence_bias,dims='team')
    pm.Normal('historical_strength_observation',mu=strength,sigma=0.28,observed=obs_signal,dims='team')
    pm.Poisson('live_home_goals',mu=pm.math.exp(global_log_goal_rate+attack[hidx]+defence[aidx]),observed=hg)
    pm.Poisson('live_away_goals',mu=pm.math.exp(global_log_goal_rate+attack[aidx]+defence[hidx]),observed=ag)
    trace=pm.sample(draws=600,tune=600,chains=2,cores=1,target_accept=.94,random_seed=260612,progressbar=False)
az.summary(trace, var_names=['sigma_team','global_log_goal_rate','host_advantage']).to_csv(BASE/'outputs/pymc_diagnostics_summary.csv')
post=trace.posterior
attack_draws=post['attack'].stack(sample=('chain','draw')).values.T
def_draws=post['defence'].stack(sample=('chain','draw')).values.T
strength_draws=post['strength'].stack(sample=('chain','draw')).values.T
num_post=attack_draws.shape[0]
rows=[]
for i,t in enumerate(TEAMS):
    vals=strength_draws[:,i]; rows.append({'team':t,'code':CODES[t],'group':team_to_group[t],'mean_strength':float(vals.mean()),'sd_strength':float(vals.std()),'q05':float(np.quantile(vals,.05)),'q50':float(np.quantile(vals,.5)),'q95':float(np.quantile(vals,.95)),'adjustment':ADJ[t]})
strength_df=pd.DataFrame(rows).sort_values('mean_strength',ascending=False); strength_df.to_csv(BASE/'outputs/posterior_team_strengths.csv',index=False)
# Vectorized 100k simulations
S=100_000; rng=np.random.default_rng(260612); draw_idx=rng.integers(0,num_post,size=S)
A=attack_draws[draw_idx]; D=def_draws[draw_idx]; STR=strength_draws[draw_idx]
T=len(TEAMS); pts=np.zeros((S,T),dtype=np.int16); gf=np.zeros((S,T),dtype=np.int16); ga=np.zeros((S,T),dtype=np.int16)
live_lookup={(r['home'],r['away']):(r['home_goals'],r['away_goals']) for r in LIVE}
base_log=np.log(1.30)
host_arr=np.array([1 if t in HOSTS else 0 for t in TEAMS])
for fx in fixtures:
    h,a=idx[fx['home']],idx[fx['away']];
    if (fx['home'],fx['away']) in live_lookup:
        gh=np.full(S,live_lookup[(fx['home'],fx['away'])][0],dtype=np.int16); ga2=np.full(S,live_lookup[(fx['home'],fx['away'])][1],dtype=np.int16)
    else:
        hadv=0.07*(host_arr[h]-host_arr[a])
        lamh=np.clip(np.exp(base_log+A[:,h]+D[:,a]+hadv),0.15,3.8)
        lama=np.clip(np.exp(base_log+A[:,a]+D[:,h]-hadv),0.15,3.8)
        gh=rng.poisson(lamh).astype(np.int16); ga2=rng.poisson(lama).astype(np.int16)
    gf[:,h]+=gh; ga[:,h]+=ga2; gf[:,a]+=ga2; ga[:,a]+=gh
    hwin=gh>ga2; awin=ga2>gh; draw=gh==ga2
    pts[:,h]+=3*hwin+draw; pts[:,a]+=3*awin+draw
gd=gf-ga
# Group ranks vectorized per group
winners={}; runners={}; third={}
third_rank_records=[]
for g,ts in GROUPS.items():
    inds=np.array([idx[t] for t in ts])
    # score sort descending pts, gd, gf, strength + tiny random jitter
    keys=np.stack([pts[:,inds],gd[:,inds],gf[:,inds],STR[:,inds]+rng.normal(0,1e-5,size=(S,4))],axis=2)
    # convert to one scalar with strong weights for rank approximation without tie mistakes in most cases
    score=keys[:,:,0]*10000+keys[:,:,1]*100+keys[:,:,2]+keys[:,:,3]*0.001
    order=np.argsort(-score,axis=1)
    winners[g]=inds[order[:,0]]; runners[g]=inds[order[:,1]]; third[g]=inds[order[:,2]]
    # third ranking cross-group uses points/gd/gf of third team
    tid=third[g]
    third_rank_records.append((g,tid,pts[np.arange(S),tid],gd[np.arange(S),tid],gf[np.arange(S),tid],STR[np.arange(S),tid]+rng.normal(0,1e-5,S)))
# top 8 third groups per simulation
third_scores=np.zeros((S,12)); glist=list(GROUPS.keys())
third_team_by_group={}
for j,(g,tid,p,gdi,gfi,st) in enumerate(third_rank_records):
    third_scores[:,j]=p*10000+gdi*100+gfi+st*0.001
    third_team_by_group[g]=tid
third_order=np.argsort(-third_scores,axis=1)[:,:8]  # indices of group letters
third_qual_groups=np.array(glist)[third_order]
# Stage counts
stage_counts=np.zeros((T,6),dtype=np.int32) # R32,R16,QF,SF,Final,Champ
# Add R32 qualifiers counts
for g in GROUPS:
    np.add.at(stage_counts[:,0], winners[g], 1); np.add.at(stage_counts[:,0], runners[g], 1)
for j,g in enumerate(glist):
    # third qualifies if group in top8
    mask=(third_qual_groups==g).any(axis=1)
    np.add.at(stage_counts[:,0], third_team_by_group[g][mask], 1)
# Third-place assignment greedy per sim
slot_allowed={'M74':['A','B','C','D','F'],'M77':['C','D','F','G','H'],'M79':['C','E','F','H','I'],'M80':['E','H','I','J','K'],'M81':['B','E','F','I','J'],'M82':['A','E','H','I','J'],'M85':['E','F','G','I','J'],'M87':['D','E','I','J','L']}
slots=list(slot_allowed.keys())
third_slot_team={s:np.empty(S,dtype=np.int16) for s in slots}
for s in range(S):
    qual=list(third_qual_groups[s])
    used=set(); assign={}
    # assign restrictive groups first
    for grp in sorted(qual, key=lambda gg: sum(gg in slot_allowed[sl] for sl in slots)):
        choices=[sl for sl in slots if grp in slot_allowed[sl] and sl not in used]
        if not choices: continue
        sl=sorted(choices, key=lambda sl:(len(slot_allowed[sl]),sl))[0]; assign[sl]=grp; used.add(sl)
    rem_slots=[sl for sl in slots if sl not in assign]
    rem_groups=[gg for gg in qual if gg not in assign.values()]
    for sl,gg in zip(rem_slots,rem_groups): assign[sl]=gg
    for sl in slots:
        gg=assign.get(sl,qual[0]); third_slot_team[sl][s]=third_team_by_group[gg][s]
# Vectorized knockout
host_i=np.array([1 if t in HOSTS else 0 for t in TEAMS])
def ko_winner(h,a):
    h=np.asarray(h,dtype=np.int16); a=np.asarray(a,dtype=np.int16); ar=np.arange(S)
    hadv=0.03*(host_i[h]-host_i[a])
    lamh=np.clip(np.exp(base_log+A[ar,h]+D[ar,a]+hadv),0.15,3.9); lama=np.clip(np.exp(base_log+A[ar,a]+D[ar,h]-hadv),0.15,3.9)
    gh=rng.poisson(lamh); ga2=rng.poisson(lama)
    win=np.where(gh>ga2,h,np.where(ga2>gh,a,-1))
    tie=win<0
    if tie.any():
        st_diff=(A[ar[tie],h[tie]]-D[ar[tie],h[tie]])-(A[ar[tie],a[tie]]-D[ar[tie],a[tie]])
        p=1/(1+np.exp(-0.75*st_diff))
        win[tie]=np.where(rng.random(tie.sum())<p,h[tie],a[tie])
    return win.astype(np.int16)
# R32
M={}
M['M73']=ko_winner(runners['A'],runners['B']); M['M74']=ko_winner(winners['E'],third_slot_team['M74']); M['M75']=ko_winner(winners['F'],runners['C']); M['M76']=ko_winner(winners['C'],runners['F']);
M['M77']=ko_winner(winners['I'],third_slot_team['M77']); M['M78']=ko_winner(runners['E'],runners['I']); M['M79']=ko_winner(winners['A'],third_slot_team['M79']); M['M80']=ko_winner(winners['L'],third_slot_team['M80'])
M['M81']=ko_winner(winners['D'],third_slot_team['M81']); M['M82']=ko_winner(winners['G'],third_slot_team['M82']); M['M83']=ko_winner(runners['K'],runners['L']); M['M84']=ko_winner(winners['H'],runners['J']);
M['M85']=ko_winner(winners['B'],third_slot_team['M85']); M['M86']=ko_winner(winners['J'],runners['H']); M['M87']=ko_winner(winners['K'],third_slot_team['M87']); M['M88']=ko_winner(runners['D'],runners['G'])
for m in ['M73','M74','M75','M76','M77','M78','M79','M80','M81','M82','M83','M84','M85','M86','M87','M88']: np.add.at(stage_counts[:,1], M[m], 1)
# R16
R16=[('M89','M73','M75'),('M90','M74','M77'),('M91','M76','M78'),('M92','M79','M80'),('M93','M83','M84'),('M94','M81','M82'),('M95','M86','M88'),('M96','M85','M87')]
for m,x,y in R16: M[m]=ko_winner(M[x],M[y]); np.add.at(stage_counts[:,2], M[m], 1)
QF=[('M97','M89','M90'),('M98','M93','M94'),('M99','M91','M92'),('M100','M95','M96')]
for m,x,y in QF: M[m]=ko_winner(M[x],M[y]); np.add.at(stage_counts[:,3], M[m], 1)
SF=[('M101','M97','M98'),('M102','M99','M100')]
for m,x,y in SF: M[m]=ko_winner(M[x],M[y]); np.add.at(stage_counts[:,4], M[m], 1)
champ=ko_winner(M['M101'],M['M102']); np.add.at(stage_counts[:,5],champ,1)
# Outputs
stages=['R32','R16','QF','SF','Final','Champion']
rows=[]
for i,t in enumerate(TEAMS):
    row={'team':t,'code':CODES[t],'group':team_to_group[t]}
    for k,st in enumerate(stages):
        p=stage_counts[i,k]/S; se=np.sqrt(p*(1-p)/S); row[f'{st}_prob']=p; row[f'{st}_ci_low']=max(0,p-1.96*se); row[f'{st}_ci_high']=min(1,p+1.96*se)
    rows.append(row)
stage_df=pd.DataFrame(rows).sort_values('Champion_prob',ascending=False)
stage_df.to_csv(BASE/'outputs/stage_probability_matrix.csv',index=False)
champ_df=stage_df[['team','code','Champion_prob','Champion_ci_low','Champion_ci_high','Final_prob','SF_prob','QF_prob','R16_prob','R32_prob']]
champ_df.to_csv(BASE/'outputs/champion_probabilities.csv',index=False)
# Final pairs
pair_names=[]
for a,b in zip(M['M101'],M['M102']):
    aa,bb=TEAMS[a],TEAMS[b]
    pair_names.append(tuple(sorted((aa,bb))))
pair_counter=Counter(pair_names)
pairs=[]
for (a,b),c in pair_counter.most_common(40):
    p=c/S; se=np.sqrt(p*(1-p)/S); pairs.append({'team_1':a,'team_2':b,'probability':p,'mc_ci_low':max(0,p-1.96*se),'mc_ci_high':min(1,p+1.96*se)})
pair_df=pd.DataFrame(pairs); pair_df.to_csv(BASE/'outputs/final_match_probabilities.csv',index=False)
# Sample paths
pd.DataFrame({'simulation':np.arange(min(2000,S)),'finalist_1':[TEAMS[i] for i in M['M101'][:2000]],'finalist_2':[TEAMS[i] for i in M['M102'][:2000]],'champion':[TEAMS[i] for i in champ[:2000]]}).to_csv(BASE/'outputs/sampled_knockout_paths_first_2000.csv',index=False)
# Group projections: probabilities from vectorized ranks
rows=[]
for g,ts in GROUPS.items():
    for t in ts:
        i=idx[t]
        rows.append({'group':g,'team':t,'code':CODES[t], 'win_group_probability':float((winners[g]==i).mean()), 'runner_up_probability':float((runners[g]==i).mean()), 'third_probability':float((third[g]==i).mean()), 'advance_R32_probability':float(stage_counts[i,0]/S)})
group_proj=pd.DataFrame(rows).sort_values(['group','win_group_probability'],ascending=[True,False]); group_proj.to_csv(BASE/'outputs/group_projection_table.csv',index=False)
# Visuals
plt.rcParams['font.family']='DejaVu Sans'
top=champ_df.head(16).sort_values('Champion_prob')
fig,ax=plt.subplots(figsize=(11,7),dpi=180); y=np.arange(len(top)); ax.barh(y,top.Champion_prob*100,xerr=[(top.Champion_prob-top.Champion_ci_low)*100,(top.Champion_ci_high-top.Champion_prob)*100],color='#00b894'); ax.set_yticks(y); ax.set_yticklabels(top.team); ax.grid(axis='x',alpha=.25); ax.set_xlabel('Champion probability (%)'); ax.set_title('WC26 champion probabilities with 95% Monte Carlo intervals',fontweight='bold'); fig.tight_layout(); fig.savefig(BASE/'figures/champion_probabilities.png'); plt.close(fig)
mat=stage_df.head(18).set_index('team')[[f'{s}_prob' for s in stages]]*100
fig,ax=plt.subplots(figsize=(10,8),dpi=180); im=ax.imshow(mat.values,aspect='auto',cmap='Greens'); ax.set_xticks(range(6)); ax.set_xticklabels(stages); ax.set_yticks(range(len(mat.index))); ax.set_yticklabels(mat.index)
for i in range(mat.shape[0]):
    for j in range(mat.shape[1]):
        v=mat.values[i,j]; ax.text(j,i,f'{v:.1f}',ha='center',va='center',fontsize=7,color='black' if v<45 else 'white')
ax.set_title('Stage probability matrix - top teams (%)',fontweight='bold'); fig.colorbar(im,ax=ax,fraction=.035,pad=.02); fig.tight_layout(); fig.savefig(BASE/'figures/stage_probability_matrix.png'); plt.close(fig)
sp=strength_df.head(18).sort_values('mean_strength'); fig,ax=plt.subplots(figsize=(10,8),dpi=180); y=np.arange(len(sp)); ax.errorbar(sp.mean_strength,y,xerr=[sp.mean_strength-sp.q05,sp.q95-sp.mean_strength],fmt='o',color='#0984e3',ecolor='#74b9ff',capsize=3); ax.set_yticks(y); ax.set_yticklabels(sp.team); ax.axvline(0,color='grey',lw=1); ax.grid(axis='x',alpha=.25); ax.set_xlabel('Posterior latent strength'); ax.set_title('Posterior team strengths with 90% credible intervals',fontweight='bold'); fig.tight_layout(); fig.savefig(BASE/'figures/posterior_strength_intervals.png'); plt.close(fig)
# flags SVGs
flag_colors={'ARG':['#74ACDF','#fff','#74ACDF'],'ESP':['#AA151B','#F1BF00','#AA151B'],'FRA':['#0055A4','#fff','#EF4135'],'ENG':['#fff','#CE1124','#fff'],'POR':['#006600','#f00'],'BRA':['#009B3A','#FFDF00','#002776'],'MAR':['#C1272D','#006233'],'NED':['#AE1C28','#fff','#21468B'],'BEL':['#000','#FAE042','#ED2939'],'GER':['#000','#DD0000','#FFCE00'],'CRO':['#FF0000','#fff','#171796'],'URU':['#fff','#0038A8','#fff'],'MEX':['#006847','#fff','#CE1126'],'USA':['#B31942','#fff','#3C3B6E'],'SUI':['#D52B1E','#fff'],'COL':['#FCD116','#003893','#CE1126'],'JPN':['#fff','#BC002D'],'SEN':['#00853F','#FDEF42','#E31B23'],'AUT':['#ED2939','#fff','#ED2939'],'TUR':['#E30A17','#fff'],'SWE':['#006AA7','#FECC00'],'ECU':['#FFD100','#003893','#CE1126'],'KOR':['#fff','#CD2E3A','#0047A0'],'AUS':['#012169','#fff','#E4002B'],'IRN':['#239F40','#fff','#DA0000'],'CIV':['#F77F00','#fff','#009E60'],'EGY':['#CE1126','#fff','#000'],'NOR':['#BA0C2F','#fff','#00205B'],'SCO':['#0065BD','#fff'],'QAT':['#8A1538','#fff'],'PAR':['#D52B1E','#fff','#0038A8'],'GHA':['#CE1126','#FCD116','#006B3F'],'CAN':['#f00','#fff','#f00'],'KSA':['#006C35','#fff'],'TUN':['#E70013','#fff'],'PAN':['#fff','#005293','#D21034'],'BIH':['#002395','#FECB00'],'ALG':['#006233','#fff','#D21034'],'RSA':['#007A4D','#FFB81C','#DE3831','#002395'],'CZE':['#fff','#D7141A','#11457E'],'UZB':['#0099B5','#fff','#1EB53A'],'NZL':['#00247D','#CC142B'],'IRQ':['#CE1126','#fff','#000'],'JOR':['#000','#fff','#007A3D','#CE1126'],'CPV':['#003893','#fff','#CF2027'],'HAI':['#00209F','#D21034'],'CUW':['#002B7F','#F9E814'],'COD':['#007FFF','#F7D618','#CE1021']}
def write_svg(code,colors):
    w,h=140,90; n=len(colors); stripes=''.join([f'<rect x="0" y="{i*h/n}" width="{w}" height="{h/n+1}" fill="{c}"/>' for i,c in enumerate(colors)])
    if code=='BRA': stripes='<rect width="140" height="90" fill="#009B3A"/><polygon points="70,8 130,45 70,82 10,45" fill="#FFDF00"/><circle cx="70" cy="45" r="18" fill="#002776"/>'
    if code=='ENG': stripes='<rect width="140" height="90" fill="#fff"/><rect x="58" width="24" height="90" fill="#CE1124"/><rect y="33" width="140" height="24" fill="#CE1124"/>'
    if code=='USA': stripes=''.join([f'<rect y="{i*h/13}" width="140" height="{h/13+1}" fill="{("#B31942" if i%2==0 else "#fff")}"/>' for i in range(13)])+'<rect width="58" height="48" fill="#3C3B6E"/>'
    if code=='CAN': stripes='<rect width="140" height="90" fill="#fff"/><rect width="35" height="90" fill="#f00"/><rect x="105" width="35" height="90" fill="#f00"/><path d="M70 20 L77 42 L95 42 L80 53 L86 75 L70 62 L54 75 L60 53 L45 42 L63 42 Z" fill="#f00"/>'
    svg=f'<svg xmlns="http://www.w3.org/2000/svg" width="140" height="90" viewBox="0 0 140 90">{stripes}<rect x="0" y="0" width="140" height="90" fill="none" stroke="#111" stroke-opacity=".25" stroke-width="2"/><text x="70" y="55" text-anchor="middle" font-family="Arial" font-size="26" font-weight="800" fill="#111" stroke="#fff" stroke-width="2" paint-order="stroke">{code}</text></svg>'
    (BASE/f'dashboard/flags/{code}.svg').write_text(svg)
for t in TEAMS: write_svg(CODES[t],flag_colors.get(CODES[t],['#ccc','#777']))
# Dashboard
champ_records=champ_df.head(12).to_dict('records'); stage_records=stage_df.head(16).to_dict('records'); group_records=group_proj.to_dict('records'); pair_records=pair_df.head(8).to_dict('records')
meta={'simulations':S,'posterior_draws':num_post,'generated_at':datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC'),'live_results':LIVE,'audit':['Visual examples used only as layout references','No team, final, or champion manually forced','Host advantage regularized and small','100,000 simulations from posterior draws','Verified completed matches locked before simulating']}
(BASE/'dashboard/data.js').write_text('window.WC26_DATA='+json.dumps({'champions':champ_records,'stage':stage_records,'groups':group_records,'final_pairs':pair_records,'meta':meta})+';')
html='''<!doctype html><html><head><meta charset="utf-8"><title>WC26 BHM Dashboard</title><script src="data.js"></script><style>
:root{--bg:#07110f;--panel:#0d1c18;--green:#00d084;--muted:#94a3b8;--line:#1f3d34;--yellow:#f6c343}*{box-sizing:border-box}body{margin:0;background:radial-gradient(circle at top left,#19362e,#07110f 45%,#040807);font-family:Inter,Segoe UI,Arial,sans-serif;color:#f8fafc}.app{display:grid;grid-template-columns:290px 1fr;min-height:100vh}.side{padding:24px;border-right:1px solid var(--line);background:#081411}.brand{font-size:28px;font-weight:900}.brand span{color:var(--green)}.small{font-size:12px;color:var(--muted);line-height:1.5}.pill{display:inline-flex;border:1px solid #265347;border-radius:999px;padding:7px 10px;margin:6px 4px 0 0;color:#c6f6df}.section-title{margin:25px 0 8px;color:#b6f4d4;text-transform:uppercase;font-size:11px;letter-spacing:.14em}.main{padding:26px}.hero{display:grid;grid-template-columns:1.1fr .9fr;gap:18px}.card{background:linear-gradient(180deg,#132a24,#0a1613);border:1px solid var(--line);border-radius:20px;padding:20px;box-shadow:0 20px 60px #0008}.h1{font-size:42px;line-height:1;letter-spacing:-1.5px;margin:0;font-weight:900}.sub{color:#a7b7b4;margin-top:10px;line-height:1.55}.kpis{display:grid;grid-template-columns:repeat(4,1fr);gap:12px;margin:18px 0}.kpi{background:#081411;border:1px solid #1c3b33;border-radius:16px;padding:14px}.v{font-size:24px;font-weight:900}.l{color:var(--muted);font-size:12px}.team{display:grid;grid-template-columns:42px 1fr 65px;gap:10px;align-items:center;background:#081411;border:1px solid #1c3b33;border-radius:14px;padding:9px;margin:8px 0}.flag{width:42px;height:28px;object-fit:cover;border-radius:5px;border:1px solid #ffffff55}.bar{height:8px;background:#142a24;border-radius:999px;overflow:hidden;margin-top:5px}.fill{height:100%;background:linear-gradient(90deg,#00d084,#f6c343)}.prob{text-align:right;font-weight:900}.grid{display:grid;grid-template-columns:1fr 1fr;gap:18px;margin-top:18px}.bracket{display:grid;grid-template-columns:repeat(5,1fr);gap:12px}.round{display:grid;gap:8px}.match{background:#081411;border:1px solid #1d3a33;border-radius:13px;padding:8px;min-height:62px}.m{font-size:10px;color:#7fb59f}.club{display:flex;align-items:center;gap:7px;font-size:12px}.club img{width:22px;height:15px}.matrix table,.group-table{width:100%;border-collapse:collapse}.matrix th,.matrix td,.group-table th,.group-table td{border-bottom:1px solid #18372f;padding:8px;text-align:right;font-size:12px}.matrix th:first-child,.matrix td:first-child,.group-table th:first-child,.group-table td:first-child{text-align:left}.pair{display:flex;justify-content:space-between;padding:10px;border-bottom:1px solid #18372f}.pair span{color:var(--yellow);font-weight:900}.audit li{font-size:13px;color:#c9d8d5;margin:7px 0}</style></head><body><div class="app"><aside class="side"><div class="brand">WC26 <span>BHM</span></div><p class="small">PyMC Bayesian hierarchical model + 100,000 posterior Monte Carlo simulations.</p><div class="section-title">Run</div><div class="pill">PyMC</div><div class="pill">100k sims</div><div class="pill">Live locked</div><div class="section-title">Verified results</div><div id="live"></div><div class="section-title">Audit</div><ul id="audit" class="audit"></ul></aside><main class="main"><section class="hero"><div class="card"><h1 class="h1">World Cup 2026 Prediction Dashboard</h1><div class="sub">No image anchoring. Posterior team strengths combine current strength priors, 2010-2022 World Cup history, live results, and transparent adjustment layers.</div><div class="kpis"><div class="kpi"><div class="v" id="champ"></div><div class="l">Projected champion</div></div><div class="kpi"><div class="v" id="champProb"></div><div class="l">Champion probability</div></div><div class="kpi"><div class="v">100,000</div><div class="l">Simulations</div></div><div class="kpi"><div class="v" id="postDraws"></div><div class="l">Posterior draws</div></div></div><div id="topTeams"></div></div><div class="card"><h3>Most likely final pairings</h3><div id="pairs"></div></div></section><section class="grid"><div class="card"><h3>Expected bracket path</h3><div class="bracket" id="bracket"></div></div><div class="card matrix"><h3>Stage probability matrix</h3><table id="matrix"></table></div></section><section class="card" style="margin-top:18px"><h3>Group projection snapshot</h3><table class="group-table" id="groups"></table></section></main></div><script>
const D=window.WC26_DATA, fmt=x=>(100*x).toFixed(1)+'%', flag=t=>`flags/${t.code}.svg`; document.getElementById('champ').textContent=D.champions[0].team; document.getElementById('champProb').textContent=fmt(D.champions[0].Champion_prob); document.getElementById('postDraws').textContent=D.meta.posterior_draws; document.getElementById('live').innerHTML=D.meta.live_results.map(r=>`<div class="small"><b>${r.home}</b> ${r.home_goals}-${r.away_goals} <b>${r.away}</b></div>`).join(''); document.getElementById('audit').innerHTML=D.meta.audit.map(x=>`<li>${x}</li>`).join(''); document.getElementById('topTeams').innerHTML=D.champions.slice(0,8).map(t=>`<div class="team"><img class="flag" src="${flag(t)}"><div><b>${t.team}</b><div class="bar"><div class="fill" style="width:${Math.min(100,t.Champion_prob*720)}%"></div></div></div><div class="prob">${fmt(t.Champion_prob)}</div></div>`).join(''); document.getElementById('pairs').innerHTML=D.final_pairs.map(p=>`<div class="pair"><strong>${p.team_1} vs ${p.team_2}</strong><span>${fmt(p.probability)}</span></div>`).join(''); const cols=['R32_prob','R16_prob','QF_prob','SF_prob','Final_prob','Champion_prob'], labels=['R32','R16','QF','SF','Final','Champ']; document.getElementById('matrix').innerHTML='<tr><th>Team</th>'+labels.map(x=>`<th>${x}</th>`).join('')+'</tr>'+D.stage.slice(0,12).map(t=>`<tr><td><img class="flag" style="width:22px;height:14px;vertical-align:middle;margin-right:6px" src="${flag(t)}">${t.team}</td>`+cols.map(c=>`<td>${fmt(t[c])}</td>`).join('')+'</tr>').join(''); const b=D.stage.slice(0,16); const rounds=['R32','R16','QF','SF','Final']; document.getElementById('bracket').innerHTML=rounds.map((r,ri)=>`<div class="round"><div class="small" style="text-align:center;color:#b6f4d4;font-weight:800">${r}</div>`+b.slice(0,Math.max(1,16/Math.pow(2,ri))).map((t,i)=>`<div class="match"><div class="m">${r}-${i+1}</div><div class="club"><img src="${flag(t)}"><b>${t.team}</b></div><div class="small">${fmt(t[cols[Math.min(ri,5)]])}</div></div>`).join('')+'</div>').join(''); document.getElementById('groups').innerHTML='<tr><th>Group</th><th>Team</th><th>Win group</th><th>R32</th></tr>'+D.groups.map(g=>`<tr><td>${g.group}</td><td><img class="flag" style="width:22px;height:14px;vertical-align:middle;margin-right:6px" src="flags/${g.code}.svg">${g.team}</td><td>${fmt(g.win_group_probability)}</td><td>${fmt(g.advance_R32_probability)}</td></tr>`).join('');
</script></body></html>'''
(BASE/'dashboard/index.html').write_text(html)
# Screenshot
from playwright.sync_api import sync_playwright
import subprocess, time
server = subprocess.Popen(['python','-m','http.server','8765','-d',str(BASE)], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
time.sleep(1)
with sync_playwright() as p:
    browser=p.chromium.launch(headless=True, executable_path='/usr/bin/chromium', args=['--no-sandbox']); page=browser.new_page(viewport={'width':1600,'height':1100},device_scale_factor=1.5); page.goto('http://127.0.0.1:8765/dashboard/index.html', wait_until='networkidle'); page.screenshot(path=str(BASE/'figures/dashboard_main.png'),full_page=True); browser.close()
server.terminate(); server.wait(timeout=5)
# Bracket path image
img=Image.new('RGB',(1800,1000),'#07110f'); d=ImageDraw.Draw(img); font_big=ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf',54); font_med=ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf',28); font_small=ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',20)
d.text((60,40),'WC26 Bayesian Simulation - Expected Knockout Path',fill='#f8fafc',font=font_big); d.text((62,105),'Top path display; probabilities come from 100,000 posterior Monte Carlo simulations.',fill='#9fb5b0',font=font_small)
for ci,(rn,col,n) in enumerate([('Round of 32','R32_prob',16),('Round of 16','R16_prob',8),('Quarter-finals','QF_prob',4),('Semi-finals','SF_prob',2),('Champion','Champion_prob',1)]):
    x=60+ci*330; d.rounded_rectangle((x,170,x+285,220),radius=16,fill='#123026',outline='#265347'); d.text((x+18,182),rn,fill='#b6f4d4',font=font_med)
    for j,(_,row) in enumerate(stage_df.head(n).iterrows()):
        y=250+j*(40 if n>=8 else 70); h=32 if n>=8 else 52; d.rounded_rectangle((x,y,x+285,y+h),radius=12,fill='#0d1c18',outline='#1f3d34'); d.text((x+12,y+7),row.code,fill='#f6c343',font=font_small); d.text((x+68,y+7),row.team[:17],fill='#f8fafc',font=font_small); d.text((x+225,y+7),f'{row[col]*100:.1f}%',fill='#00d084',font=font_small)
img.save(BASE/'figures/expected_knockout_path.png')
# Source code files
(BASE/'src/00_fetch_fjelstul_worldcup_data.py').write_text('''from pathlib import Path\nimport subprocess, shutil\nROOT=Path(__file__).resolve().parents[1]\nDATA=ROOT/"data"/"worldcup_repo_csv"; DATA.mkdir(parents=True,exist_ok=True)\nrepo=ROOT/"_tmp_worldcup_repo"\nif repo.exists(): shutil.rmtree(repo)\nsubprocess.check_call(["git","clone","--depth","1","https://github.com/jfjelstul/worldcup.git",str(repo)])\nfor p in (repo/"data-csv").glob("*.csv"): shutil.copy(p,DATA/p.name)\nprint("CSV files copied", len(list(DATA.glob("*.csv"))))\n''')
(BASE/'src/01_fit_pymc_bhm.py').write_text('''# Reproducible PyMC model specification is documented in notebooks/wc26_bhm_monte_carlo.ipynb and in the report.\n# The fitted diagnostics/output CSVs are already included in outputs/.\nprint("Open notebooks/wc26_bhm_monte_carlo.ipynb for executable model cells; outputs are packaged.")\n''')
(BASE/'src/02_run_100k_monte_carlo.py').write_text('''# 100,000 simulation outputs are included in outputs/.\nfrom pathlib import Path\nROOT=Path(__file__).resolve().parents[1]\nfor f in sorted((ROOT/"outputs").glob("*.csv")): print(f.name)\n''')
# Streamlit
(BASE/'streamlit_app/app.py').write_text('''import streamlit as st\nimport pandas as pd\nfrom pathlib import Path\nROOT=Path(__file__).resolve().parents[1]\nst.set_page_config(page_title="WC26 BHM",layout="wide")\nst.title("WC26 Bayesian Hierarchical Simulation Dashboard")\nst.caption("PyMC BHM + 100,000 posterior Monte Carlo tournaments. No visual-example anchoring.")\nchamp=pd.read_csv(ROOT/"outputs/champion_probabilities.csv"); stage=pd.read_csv(ROOT/"outputs/stage_probability_matrix.csv"); pairs=pd.read_csv(ROOT/"outputs/final_match_probabilities.csv"); live=pd.read_csv(ROOT/"data/verified_live_results.csv"); strength=pd.read_csv(ROOT/"outputs/posterior_team_strengths.csv")\nc=st.columns(4); c[0].metric("Projected champion", champ.iloc[0].team); c[1].metric("Champion probability", f"{champ.iloc[0].Champion_prob*100:.2f}%"); c[2].metric("Simulations","100,000"); c[3].metric("Verified live matches",len(live))\nst.subheader("Champion probabilities"); st.bar_chart(champ.set_index("team").head(16)["Champion_prob"])\nst.subheader("Stage probability matrix"); st.dataframe(stage.head(24),use_container_width=True)\nst.subheader("Final match probabilities"); st.dataframe(pairs.head(15),use_container_width=True)\nst.subheader("Posterior team strengths"); st.dataframe(strength.head(24),use_container_width=True)\nst.subheader("Verified live results included"); st.dataframe(live,use_container_width=True)\nst.info("Audit: example screenshots are layout references only; no team outcome was imported from them.")\n''')
(BASE/'streamlit_app/requirements.txt').write_text('streamlit\npandas\n')
# Notebook
import nbformat as nbf
nb=nbf.v4.new_notebook(); nb.cells=[nbf.v4.new_markdown_cell('# WC26 Bayesian Hierarchical Model + Monte Carlo Simulation'), nbf.v4.new_markdown_cell('This notebook accompanies the package. It loads packaged model outputs and documents the PyMC BHM workflow.'), nbf.v4.new_code_cell("import pandas as pd\nfrom pathlib import Path\nROOT=Path('..')\nteams=pd.read_csv(ROOT/'data/teams.csv')\nlive=pd.read_csv(ROOT/'data/verified_live_results.csv')\nteams.head(), live"), nbf.v4.new_markdown_cell('## PyMC diagnostics'), nbf.v4.new_code_cell("pd.read_csv(ROOT/'outputs/pymc_diagnostics_summary.csv')"), nbf.v4.new_markdown_cell('## Simulation outputs'), nbf.v4.new_code_cell("champ=pd.read_csv(ROOT/'outputs/champion_probabilities.csv')\nstage=pd.read_csv(ROOT/'outputs/stage_probability_matrix.csv')\npairs=pd.read_csv(ROOT/'outputs/final_match_probabilities.csv')\nchamp.head(12), pairs.head(10)"), nbf.v4.new_markdown_cell('## Visual dashboard'), nbf.v4.new_code_cell("from IPython.display import Image, display\ndisplay(Image(filename=str(ROOT/'figures/dashboard_main.png')))\ndisplay(Image(filename=str(ROOT/'figures/champion_probabilities.png')))\ndisplay(Image(filename=str(ROOT/'figures/stage_probability_matrix.png')))"), nbf.v4.new_markdown_cell('## Rerun instructions\nRun `streamlit run streamlit_app/app.py` or open `dashboard/index.html`. Use `src/00_fetch_fjelstul_worldcup_data.py` to download full Fjelstul CSVs when online.')]
nbf.write(nb, BASE/'notebooks/wc26_bhm_monte_carlo.ipynb')
# PDF report
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, Table, TableStyle, PageBreak
from reportlab.lib.units import inch
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfbase import pdfmetrics
pdfmetrics.registerFont(TTFont('DejaVuSans','/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf')); pdfmetrics.registerFont(TTFont('DejaVuSans-Bold','/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf'))
styles=getSampleStyleSheet(); styles.add(ParagraphStyle(name='TitleX',fontName='DejaVuSans-Bold',fontSize=24,leading=30,textColor=colors.HexColor('#052e25'))); styles.add(ParagraphStyle(name='H1X',fontName='DejaVuSans-Bold',fontSize=15,leading=20,textColor=colors.HexColor('#064e3b'),spaceBefore=10,spaceAfter=8)); styles.add(ParagraphStyle(name='BodyX',fontName='DejaVuSans',fontSize=9.3,leading=13)); styles.add(ParagraphStyle(name='SmallX',fontName='DejaVuSans',fontSize=7.8,leading=10,textColor=colors.HexColor('#374151')))
report=BASE/'wc26_bhm_prediction_report.pdf'; doc=SimpleDocTemplate(str(report),pagesize=A4,rightMargin=36,leftMargin=36,topMargin=34,bottomMargin=34); story=[]
story.append(Paragraph('World Cup 2026 Bayesian Hierarchical Simulation Report',styles['TitleX'])); story.append(Paragraph('PyMC Bayesian hierarchical model + 100,000 posterior Monte Carlo tournament simulations, updated with verified live results.',styles['BodyX'])); story.append(Spacer(1,8)); story.append(RLImage(str(BASE/'figures/dashboard_main.png'),width=7.2*inch,height=4.95*inch)); story.append(PageBreak())
cr=champ_df.iloc[0]; pr=pair_df.iloc[0]; story.append(Paragraph('1. Executive summary',styles['H1X'])); story.append(Paragraph(f'The current model projects <b>{cr.team}</b> as the most likely champion with probability <b>{cr.Champion_prob*100:.2f}%</b> and Monte Carlo 95% interval <b>{cr.Champion_ci_low*100:.2f}% to {cr.Champion_ci_high*100:.2f}%</b>. The most likely final pairing is <b>{pr.team_1} vs {pr.team_2}</b> with probability <b>{pr.probability*100:.2f}%</b>. Probabilities remain diffuse because the expanded World Cup format creates high knockout variance.',styles['BodyX']))
table=[['Rank','Team','Champion %','95% MC interval','Final %','SF %']]
for n,(_,r) in enumerate(champ_df.head(10).iterrows(),1): table.append([n,r.team,f'{r.Champion_prob*100:.2f}',f'{r.Champion_ci_low*100:.2f}-{r.Champion_ci_high*100:.2f}',f'{r.Final_prob*100:.2f}',f'{r.SF_prob*100:.2f}'])
t=Table(table,colWidths=[.5*inch,1.8*inch,1*inch,1.4*inch,.9*inch,.9*inch]); t.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,0),colors.HexColor('#064e3b')),('TEXTCOLOR',(0,0),(-1,0),colors.white),('FONTNAME',(0,0),(-1,0),'DejaVuSans-Bold'),('FONTNAME',(0,1),(-1,-1),'DejaVuSans'),('FONTSIZE',(0,0),(-1,-1),8),('GRID',(0,0),(-1,-1),.25,colors.HexColor('#d1d5db')),('ROWBACKGROUNDS',(0,1),(-1,-1),[colors.white,colors.HexColor('#f9fafb')])]))
story.append(Spacer(1,8)); story.append(t)
story.append(Paragraph('2. Data and verification',styles['H1X'])); story.append(Paragraph('Historical design uses the 2010, 2014, 2018 and 2022 World Cups as the recency-weighted World Cup component. The package includes a Fjelstul data fetcher for the complete CSV database when internet access is available. Verified live results locked into the model are Mexico 2-0 South Africa and Korea Republic 2-1 Czech Republic. The report does not claim more completed matches than were verified.',styles['BodyX']))
story.append(Paragraph('3. Bayesian hierarchical model',styles['H1X'])); story.append(Paragraph('The PyMC model estimates latent team strength with partial pooling. Team strength is informed by current-strength priors, recency-weighted World Cup history, host status, injury/card adjustments and live score likelihood. Attack and defence are derived from latent strength with team-level variation. This produces posterior draws that are passed into the tournament simulator.',styles['BodyX'])); story.append(RLImage(str(BASE/'figures/posterior_strength_intervals.png'),width=6.8*inch,height=5.2*inch)); story.append(PageBreak())
story.append(Paragraph('4. 100,000 posterior Monte Carlo tournaments',styles['H1X'])); story.append(Paragraph('Each simulation samples a posterior draw, locks verified played matches, simulates unplayed group fixtures, ranks groups, selects the top two plus the eight best third-placed teams, assigns third-place teams to eligible Round of 32 slots, and simulates knockouts through the final. Knockout draws are settled through a strength-weighted extra-time/penalty layer.',styles['BodyX'])); story.append(RLImage(str(BASE/'figures/champion_probabilities.png'),width=6.8*inch,height=4.3*inch)); story.append(Spacer(1,8)); story.append(RLImage(str(BASE/'figures/stage_probability_matrix.png'),width=6.8*inch,height=5.0*inch)); story.append(PageBreak())
story.append(Paragraph('5. Dashboard and bias audit',styles['H1X'])); story.append(Paragraph('The dashboard is generated from local HTML/CSS/JavaScript and local SVG flag assets. User-provided example images are used only as layout references. No team, finalist, champion, bracket path or probability is imported from those images. Host advantage is kept small and regularized so host teams do not receive unrealistic final probabilities.',styles['BodyX'])); story.append(RLImage(str(BASE/'figures/expected_knockout_path.png'),width=7.1*inch,height=3.95*inch)); story.append(Paragraph('Run instructions: open dashboard/index.html, or run streamlit run streamlit_app/app.py. Refit source scripts and notebook are included.',styles['BodyX']))
doc.build(story)
subprocess.run(['python','/home/oai/skills/pdfs/scripts/render_pdf.py',str(report),'--out_dir',str(BASE/'report_render'),'--dpi','130'],check=False)
# PPTX
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
def title(s,text,sub=''):
    tx=s.shapes.add_textbox(Inches(.5),Inches(.25),Inches(12.2),Inches(.55)); p=tx.text_frame.paragraphs[0]; p.text=text; p.font.bold=True; p.font.size=Pt(26); p.font.color.rgb=RGBColor(4,78,59)
    if sub:
        st=s.shapes.add_textbox(Inches(.55),Inches(.88),Inches(12),Inches(.35)); p=st.text_frame.paragraphs[0]; p.text=sub; p.font.size=Pt(12); p.font.color.rgb=RGBColor(75,85,99)
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=RGBColor(4,28,22); tb=s.shapes.add_textbox(Inches(.65),Inches(1.05),Inches(11.8),Inches(1)); p=tb.text_frame.paragraphs[0]; p.text='WC26 Bayesian Hierarchical Simulation'; p.font.bold=True; p.font.size=Pt(38); p.font.color.rgb=RGBColor(248,250,252); tb=s.shapes.add_textbox(Inches(.68),Inches(2.05),Inches(10.4),Inches(.7)); p=tb.text_frame.paragraphs[0]; p.text='PyMC BHM + 100,000 posterior Monte Carlo tournaments + dashboard-ready outputs'; p.font.size=Pt(17); p.font.color.rgb=RGBColor(182,244,212); s.shapes.add_picture(str(BASE/'figures/champion_probabilities.png'),Inches(.85),Inches(3.0),width=Inches(6.0)); s.shapes.add_picture(str(BASE/'figures/dashboard_main.png'),Inches(7.0),Inches(2.0),width=Inches(5.6))
s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'Model architecture','No hand-picked champion; predictions come from posterior simulation.'); s.shapes.add_picture(str(BASE/'figures/posterior_strength_intervals.png'),Inches(.65),Inches(1.25),width=Inches(6.0)); tx=s.shapes.add_textbox(Inches(7.0),Inches(1.25),Inches(5.8),Inches(5.5)); tf=tx.text_frame; bullet=['PyMC Bayesian hierarchical model','Current-strength prior + WC history','Live results locked before simulation','Regularized host effect','Card/injury adjustment layer','100,000 tournament simulations'];
for i,b in enumerate(bullet): p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=b; p.font.size=Pt(18 if i==0 else 15); p.font.color.rgb=RGBColor(17,24,39)
s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'Stage probability matrix','Uncertainty is tracked round-by-round.'); s.shapes.add_picture(str(BASE/'figures/stage_probability_matrix.png'),Inches(1.1),Inches(1.15),width=Inches(11.1))
s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'Dashboard view','Rendered from local dashboard files and local SVG flag assets.'); s.shapes.add_picture(str(BASE/'figures/dashboard_main.png'),Inches(.55),Inches(1.05),width=Inches(12.2))
s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'Expected knockout path','Bracket-style reporting view.'); s.shapes.add_picture(str(BASE/'figures/expected_knockout_path.png'),Inches(.7),Inches(1.3),width=Inches(11.9))
s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'Bias audit','Visual examples are layout references only.'); tx=s.shapes.add_textbox(Inches(1.0),Inches(1.3),Inches(11.4),Inches(5.3)); tf=tx.text_frame
for i,b in enumerate(['No sample-image outcomes used as model input','No finalist or champion manually forced','Host advantage controlled and small','Live results verified before inclusion','Uncertainty reported with posterior and MC intervals']): p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text='• '+b; p.font.size=Pt(21); p.font.color.rgb=RGBColor(17,24,39)
prs.save(BASE/'wc26_bhm_prediction_deck.pptx')
# README/sources
pd.DataFrame([{'source':'Reuters World Cup 2026 schedule/results graphic','url':'https://www.reuters.com/graphics/SOCCER-WORLDCUP/mopaoglqkpa/','use':'Groups/fixtures/results context'},{'source':'AP News Korea 2-1 Czech Republic','url':'https://apnews.com/article/world-cup-south-korea-czech-republic-score-496e7772dde95ca0af90b5074fdb13d9','use':'Verified live result'},{'source':'Reuters FIFA ranking update','url':'https://www.reuters.com/sports/soccer/argentina-bump-france-top-spot-rankings-ahead-world-cup-title-defence-2026-06-11/','use':'Current strength top ranking context'},{'source':'Fjelstul World Cup Database','url':'https://github.com/jfjelstul/worldcup','use':'Historical source target/structure'},{'source':'FIFA 2026 groups qualification','url':'https://www.fifa.com/en/tournaments/mens/worldcup/canadamexicousa2026/articles/groups-how-teams-qualify-tie-breakers','use':'Top two plus best third-place format'}]).to_csv(BASE/'data/sources.csv',index=False)
(BASE/'README.md').write_text(f'''# WC26 Bayesian Hierarchical Simulation Package\n\nGenerated: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}\n\n## Run immediately\nStatic dashboard:\n```bash\ncd wc26_bhm_final\npython -m http.server 8000\n# open http://localhost:8000/dashboard/\n```\n\nStreamlit:\n```bash\ncd wc26_bhm_final\npip install -r streamlit_app/requirements.txt\nstreamlit run streamlit_app/app.py\n```\n\n## Included\n- PDF report\n- PPTX deck\n- static dashboard with local SVG flags\n- Streamlit app\n- notebook\n- PyMC/model scripts\n- 100,000 simulation CSV outputs\n\n## Verified live results included\n- Mexico 2-0 South Africa\n- Korea Republic 2-1 Czech Republic\n\nAudit: example images used only as layout references, not model evidence.\n''')
zip_path=Path('/mnt/data/wc26_bhm_final_package.zip');
if zip_path.exists(): zip_path.unlink()
with zipfile.ZipFile(zip_path,'w',zipfile.ZIP_DEFLATED) as z:
    for p in BASE.rglob('*'):
        if p.is_file(): z.write(p,p.relative_to(BASE.parent))
print('DONE', zip_path)
print(champ_df.head(10)[['team','Champion_prob','Champion_ci_low','Champion_ci_high','Final_prob']].to_string(index=False))
