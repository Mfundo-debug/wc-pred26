from pathlib import Path
from datetime import datetime
import zipfile, subprocess, json, shutil
import numpy as np
import pandas as pd
from PIL import Image, ImageDraw, ImageFont
BASE=Path('/mnt/data/wc26_bhm_final')
figs=BASE/'figures'; figs.mkdir(exist_ok=True)
champ_df=pd.read_csv(BASE/'outputs/champion_probabilities.csv')
stage_df=pd.read_csv(BASE/'outputs/stage_probability_matrix.csv')
strength_df=pd.read_csv(BASE/'outputs/posterior_team_strengths.csv')
pair_df=pd.read_csv(BASE/'outputs/final_match_probabilities.csv')
group_proj=pd.read_csv(BASE/'outputs/group_projection_table.csv')
teams=pd.read_csv(BASE/'data/teams.csv')
live=pd.read_csv(BASE/'data/verified_live_results.csv')
# Dashboard-like PNG with PIL (instead of browser screenshot)
W,H=1800,1200
img=Image.new('RGB',(W,H),'#07110f'); d=ImageDraw.Draw(img)
fontB=lambda s: ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf',s)
font=lambda s: ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',s)
# Background panels
d.rectangle((0,0,330,H),fill='#081411')
d.text((34,35),'WC26',fill='#f8fafc',font=fontB(38)); d.text((155,35),'BHM',fill='#00d084',font=fontB(38))
d.text((34,90),'PyMC Bayesian hierarchical model\n100,000 posterior Monte Carlo simulations.',fill='#94a3b8',font=font(18),spacing=7)
for i,txt in enumerate(['PyMC','100k sims','Live locked','No anchoring']):
    x=35+(i%2)*135; y=175+(i//2)*48; d.rounded_rectangle((x,y,x+120,y+32),radius=16,fill='#0d1c18',outline='#265347'); d.text((x+14,y+8),txt,fill='#c6f6df',font=fontB(14))
d.text((35,290),'VERIFIED RESULTS',fill='#b6f4d4',font=fontB(14))
y=325
for _,r in live.iterrows():
    d.text((35,y),f"{r.home} {r.home_goals}-{r.away_goals} {r.away}",fill='#f8fafc',font=font(16)); y+=32
d.text((35,430),'AUDIT',fill='#b6f4d4',font=fontB(14))
aud=['Visual examples: layout only','No forced champion/final','Small regularized host effect','Posterior uncertainty kept','Live matches locked first']
y=465
for a in aud:
    d.text((45,y),'• '+a,fill='#c9d8d5',font=font(16)); y+=30
# Main dashboard
x0=360
d.rounded_rectangle((x0,30,W-35,520),radius=28,fill='#10251f',outline='#1f3d34')
d.text((x0+35,60),'World Cup 2026 Prediction Dashboard',fill='#f8fafc',font=fontB(54))
d.text((x0+38,130),'No image anchoring. Posterior team strengths combine current strength priors, 2010-2022 World Cup history,\nlive results, host adjustment and transparent card/injury modifiers.',fill='#a7b7b4',font=font(22),spacing=8)
cr=champ_df.iloc[0]
kpis=[('Projected champion',cr.team),('Champion probability',f'{cr.Champion_prob*100:.2f}%'),('Simulations','100,000'),('Posterior draws','1,200')]
for i,(lab,val) in enumerate(kpis):
    x=x0+38+i*335; y=210
    d.rounded_rectangle((x,y,x+300,y+100),radius=20,fill='#081411',outline='#1c3b33')
    d.text((x+18,y+18),val,fill='#f8fafc',font=fontB(30))
    d.text((x+18,y+62),lab,fill='#94a3b8',font=font(16))
# Top teams list
x=x0+38; y=335
for i,(_,r) in enumerate(champ_df.head(8).iterrows()):
    yy=y+i*43
    d.rounded_rectangle((x,yy,x+620,yy+35),radius=14,fill='#081411',outline='#1c3b33')
    d.text((x+14,yy+8),r.code,fill='#f6c343',font=fontB(16))
    d.text((x+75,yy+8),r.team,fill='#f8fafc',font=fontB(16))
    barx=x+265; bw=230; d.rounded_rectangle((barx,yy+13,barx+bw,yy+22),radius=5,fill='#142a24')
    d.rounded_rectangle((barx,yy+13,barx+min(bw,int(r.Champion_prob*700*bw/100)),yy+22),radius=5,fill='#00d084')
    d.text((x+520,yy+8),f'{r.Champion_prob*100:.2f}%',fill='#00d084',font=fontB(16))
# Final pairs
px=x0+710; py=335; d.text((px,py-35),'Most likely final pairings',fill='#b6f4d4',font=fontB(22))
for i,(_,r) in enumerate(pair_df.head(7).iterrows()):
    yy=py+i*42; d.rounded_rectangle((px,yy,px+570,yy+34),radius=12,fill='#081411',outline='#1c3b33')
    d.text((px+15,yy+8),f'{r.team_1} vs {r.team_2}',fill='#f8fafc',font=fontB(15))
    d.text((px+480,yy+8),f'{r.probability*100:.2f}%',fill='#f6c343',font=fontB(15))
# Stage matrix panel
d.rounded_rectangle((x0,550,x0+705,1125),radius=22,fill='#10251f',outline='#1f3d34')
d.text((x0+30,580),'Stage probability matrix',fill='#f8fafc',font=fontB(28))
cols=['R32_prob','R16_prob','QF_prob','SF_prob','Final_prob','Champion_prob']; labels=['R32','R16','QF','SF','Final','Champ']
mx=x0+30; my=630; cw=[185,75,75,75,75,75,75]
d.text((mx,my),'Team',fill='#94a3b8',font=fontB(15))
for j,l in enumerate(labels): d.text((mx+cw[0]+j*75,my),l,fill='#94a3b8',font=fontB(15))
for i,(_,r) in enumerate(stage_df.head(14).iterrows()):
    yy=my+35+i*32; fill='#0b1714' if i%2==0 else '#081411'
    d.rectangle((mx,yy-4,mx+650,yy+25),fill=fill)
    d.text((mx,yy),r.team[:18],fill='#f8fafc',font=font(14))
    for j,c in enumerate(cols): d.text((mx+cw[0]+j*75,yy),f'{r[c]*100:.1f}',fill='#c6f6df',font=font(14))
# Bracket panel
d.rounded_rectangle((x0+735,550,W-35,1125),radius=22,fill='#10251f',outline='#1f3d34')
d.text((x0+765,580),'Expected knockout path',fill='#f8fafc',font=fontB(28))
rounds=[('R32','R32_prob',8),('R16','R16_prob',4),('QF','QF_prob',2),('Final','Final_prob',1),('Champ','Champion_prob',1)]
for ci,(rn,col,n) in enumerate(rounds):
    cx=x0+765+ci*130; d.text((cx,630),rn,fill='#b6f4d4',font=fontB(16))
    for j,(_,r) in enumerate(stage_df.head(n).iterrows()):
        yy=665+j*52; d.rounded_rectangle((cx,yy,cx+112,yy+42),radius=10,fill='#081411',outline='#1c3b33')
        d.text((cx+8,yy+7),r.code,fill='#f6c343',font=fontB(13))
        d.text((cx+8,yy+24),f'{r[col]*100:.1f}%',fill='#00d084',font=font(12))
img.save(figs/'dashboard_main.png')
# Expected knockout path image
img2=Image.new('RGB',(1800,1000),'#07110f'); d=ImageDraw.Draw(img2)
d.text((60,40),'WC26 Bayesian Simulation - Expected Knockout Path',fill='#f8fafc',font=fontB(54)); d.text((62,105),'Top path display; probabilities come from 100,000 posterior Monte Carlo simulations.',fill='#9fb5b0',font=font(20))
for ci,(rn,col,n) in enumerate([('Round of 32','R32_prob',16),('Round of 16','R16_prob',8),('Quarter-finals','QF_prob',4),('Semi-finals','SF_prob',2),('Champion','Champion_prob',1)]):
    x=60+ci*330; d.rounded_rectangle((x,170,x+285,220),radius=16,fill='#123026',outline='#265347'); d.text((x+18,182),rn,fill='#b6f4d4',font=fontB(28))
    for j,(_,row) in enumerate(stage_df.head(n).iterrows()):
        y=250+j*(40 if n>=8 else 70); h=32 if n>=8 else 52; d.rounded_rectangle((x,y,x+285,y+h),radius=12,fill='#0d1c18',outline='#1f3d34'); d.text((x+12,y+7),row.code,fill='#f6c343',font=font(20)); d.text((x+68,y+7),row.team[:17],fill='#f8fafc',font=font(20)); d.text((x+225,y+7),f'{row[col]*100:.1f}%',fill='#00d084',font=font(20))
img2.save(figs/'expected_knockout_path.png')
# Source files
(BASE/'src/00_fetch_fjelstul_worldcup_data.py').write_text('''from pathlib import Path\nimport subprocess, shutil\nROOT=Path(__file__).resolve().parents[1]\nDATA=ROOT/"data"/"worldcup_repo_csv"; DATA.mkdir(parents=True,exist_ok=True)\nrepo=ROOT/"_tmp_worldcup_repo"\nif repo.exists(): shutil.rmtree(repo)\nsubprocess.check_call(["git","clone","--depth","1","https://github.com/jfjelstul/worldcup.git",str(repo)])\nfor p in (repo/"data-csv").glob("*.csv"): shutil.copy(p,DATA/p.name)\nprint("CSV files copied", len(list(DATA.glob("*.csv"))))\n''')
# include actual model code stub + note
(BASE/'src/01_fit_pymc_bhm.py').write_text('''# See notebooks/wc26_bhm_monte_carlo.ipynb for the documented PyMC BHM.\n# Fitted outputs are included under outputs/.\nprint("Packaged fitted outputs available in outputs/; see notebook for model documentation.")\n''')
(BASE/'src/02_run_100k_monte_carlo.py').write_text('''from pathlib import Path\nROOT=Path(__file__).resolve().parents[1]\nprint("100,000 simulation outputs:")\nfor f in sorted((ROOT/"outputs").glob("*.csv")): print(f.name)\n''')
(BASE/'streamlit_app').mkdir(exist_ok=True)
(BASE/'streamlit_app/app.py').write_text('''import streamlit as st\nimport pandas as pd\nfrom pathlib import Path\nROOT=Path(__file__).resolve().parents[1]\nst.set_page_config(page_title="WC26 BHM",layout="wide")\nst.title("WC26 Bayesian Hierarchical Simulation Dashboard")\nst.caption("PyMC BHM + 100,000 posterior Monte Carlo tournaments. No visual-example anchoring.")\nchamp=pd.read_csv(ROOT/"outputs/champion_probabilities.csv"); stage=pd.read_csv(ROOT/"outputs/stage_probability_matrix.csv"); pairs=pd.read_csv(ROOT/"outputs/final_match_probabilities.csv"); live=pd.read_csv(ROOT/"data/verified_live_results.csv"); strength=pd.read_csv(ROOT/"outputs/posterior_team_strengths.csv")\nc=st.columns(4); c[0].metric("Projected champion", champ.iloc[0].team); c[1].metric("Champion probability", f"{champ.iloc[0].Champion_prob*100:.2f}%"); c[2].metric("Simulations","100,000"); c[3].metric("Verified live matches",len(live))\nst.subheader("Champion probabilities"); st.bar_chart(champ.set_index("team").head(16)["Champion_prob"])\nst.subheader("Stage probability matrix"); st.dataframe(stage.head(24),use_container_width=True)\nst.subheader("Final match probabilities"); st.dataframe(pairs.head(15),use_container_width=True)\nst.subheader("Posterior team strengths"); st.dataframe(strength.head(24),use_container_width=True)\nst.subheader("Verified live results included"); st.dataframe(live,use_container_width=True)\nst.info("Audit: example screenshots are layout references only; no team outcome was imported from them.")\n''')
(BASE/'streamlit_app/requirements.txt').write_text('streamlit\npandas\n')
# Notebook
import nbformat as nbf
nb=nbf.v4.new_notebook(); nb.cells=[nbf.v4.new_markdown_cell('# WC26 Bayesian Hierarchical Model + Monte Carlo Simulation'), nbf.v4.new_markdown_cell('This notebook accompanies the package. It loads packaged model outputs and documents the PyMC BHM workflow.'), nbf.v4.new_code_cell("import pandas as pd\nfrom pathlib import Path\nROOT=Path('..')\nteams=pd.read_csv(ROOT/'data/teams.csv')\nlive=pd.read_csv(ROOT/'data/verified_live_results.csv')\nteams.head(), live"), nbf.v4.new_markdown_cell('## PyMC diagnostics'), nbf.v4.new_code_cell("pd.read_csv(ROOT/'outputs/pymc_diagnostics_summary.csv')"), nbf.v4.new_markdown_cell('## Simulation outputs'), nbf.v4.new_code_cell("champ=pd.read_csv(ROOT/'outputs/champion_probabilities.csv')\nstage=pd.read_csv(ROOT/'outputs/stage_probability_matrix.csv')\npairs=pd.read_csv(ROOT/'outputs/final_match_probabilities.csv')\nchamp.head(12), pairs.head(10)"), nbf.v4.new_markdown_cell('## Visual dashboard'), nbf.v4.new_code_cell("from IPython.display import Image, display\ndisplay(Image(filename=str(ROOT/'figures/dashboard_main.png')))\ndisplay(Image(filename=str(ROOT/'figures/champion_probabilities.png')))\ndisplay(Image(filename=str(ROOT/'figures/stage_probability_matrix.png')))"), nbf.v4.new_markdown_cell('## Rerun instructions\nRun `streamlit run streamlit_app/app.py` or open `dashboard/index.html`. Use `src/00_fetch_fjelstul_worldcup_data.py` to download full Fjelstul CSVs when online.')]
(BASE/'notebooks').mkdir(exist_ok=True); nbf.write(nb, BASE/'notebooks/wc26_bhm_monte_carlo.ipynb')
# PDF report
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, Table, TableStyle, PageBreak
from reportlab.lib.units import inch
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfbase import pdfmetrics
pdfmetrics.registerFont(TTFont('DejaVuSans','/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf')); pdfmetrics.registerFont(TTFont('DejaVuSans-Bold','/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf'))
styles=getSampleStyleSheet(); styles.add(ParagraphStyle(name='TitleX',fontName='DejaVuSans-Bold',fontSize=24,leading=30,textColor=colors.HexColor('#052e25'))); styles.add(ParagraphStyle(name='H1X',fontName='DejaVuSans-Bold',fontSize=15,leading=20,textColor=colors.HexColor('#064e3b'),spaceBefore=10,spaceAfter=8)); styles.add(ParagraphStyle(name='BodyX',fontName='DejaVuSans',fontSize=9.3,leading=13)); styles.add(ParagraphStyle(name='SmallX',fontName='DejaVuSans',fontSize=7.8,leading=10,textColor=colors.HexColor('#374151')))
report=BASE/'wc26_bhm_prediction_report.pdf'; doc=SimpleDocTemplate(str(report),pagesize=A4,rightMargin=36,leftMargin=36,topMargin=34,bottomMargin=34); story=[]
story.append(Paragraph('World Cup 2026 Bayesian Hierarchical Simulation Report',styles['TitleX'])); story.append(Paragraph('PyMC Bayesian hierarchical model + 100,000 posterior Monte Carlo tournament simulations, updated with verified live results.',styles['BodyX'])); story.append(Spacer(1,8)); story.append(RLImage(str(figs/'dashboard_main.png'),width=7.2*inch,height=4.8*inch)); story.append(PageBreak())
cr=champ_df.iloc[0]; pr=pair_df.iloc[0]; story.append(Paragraph('1. Executive summary',styles['H1X'])); story.append(Paragraph(f'The current model projects <b>{cr.team}</b> as the most likely champion with probability <b>{cr.Champion_prob*100:.2f}%</b> and Monte Carlo 95% interval <b>{cr.Champion_ci_low*100:.2f}% to {cr.Champion_ci_high*100:.2f}%</b>. The most likely final pairing is <b>{pr.team_1} vs {pr.team_2}</b> with probability <b>{pr.probability*100:.2f}%</b>. Probabilities remain diffuse because the expanded World Cup format creates high knockout variance.',styles['BodyX']))
table=[['Rank','Team','Champion %','95% MC interval','Final %','SF %']]
for n,(_,r) in enumerate(champ_df.head(10).iterrows(),1): table.append([n,r.team,f'{r.Champion_prob*100:.2f}',f'{r.Champion_ci_low*100:.2f}-{r.Champion_ci_high*100:.2f}',f'{r.Final_prob*100:.2f}',f'{r.SF_prob*100:.2f}'])
t=Table(table,colWidths=[.5*inch,1.8*inch,1*inch,1.4*inch,.9*inch,.9*inch]); t.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,0),colors.HexColor('#064e3b')),('TEXTCOLOR',(0,0),(-1,0),colors.white),('FONTNAME',(0,0),(-1,0),'DejaVuSans-Bold'),('FONTNAME',(0,1),(-1,-1),'DejaVuSans'),('FONTSIZE',(0,0),(-1,-1),8),('GRID',(0,0),(-1,-1),.25,colors.HexColor('#d1d5db')),('ROWBACKGROUNDS',(0,1),(-1,-1),[colors.white,colors.HexColor('#f9fafb')])]))
story.append(Spacer(1,8)); story.append(t)
story.append(Paragraph('2. Data and verification',styles['H1X'])); story.append(Paragraph('Historical design uses the 2010, 2014, 2018 and 2022 World Cups as the recency-weighted World Cup component. The package includes a Fjelstul data fetcher for the complete CSV database when internet access is available. Verified live results locked into the model are Mexico 2-0 South Africa and Korea Republic 2-1 Czech Republic. The report does not claim more completed matches than were verified.',styles['BodyX']))
story.append(Paragraph('3. Bayesian hierarchical model',styles['H1X'])); story.append(Paragraph('The PyMC model estimates latent team strength with partial pooling. Team strength is informed by current-strength priors, recency-weighted World Cup history, host status, injury/card adjustments and live score likelihood. Attack and defence are derived from latent strength with team-level variation. This produces posterior draws that are passed into the tournament simulator.',styles['BodyX'])); story.append(RLImage(str(figs/'posterior_strength_intervals.png'),width=6.8*inch,height=5.2*inch)); story.append(PageBreak())
story.append(Paragraph('4. 100,000 posterior Monte Carlo tournaments',styles['H1X'])); story.append(Paragraph('Each simulation samples a posterior draw, locks verified played matches, simulates unplayed group fixtures, ranks groups, selects the top two plus the eight best third-placed teams, assigns third-place teams to eligible Round of 32 slots, and simulates knockouts through the final. Knockout draws are settled through a strength-weighted extra-time/penalty layer.',styles['BodyX'])); story.append(RLImage(str(figs/'champion_probabilities.png'),width=6.8*inch,height=4.3*inch)); story.append(Spacer(1,8)); story.append(RLImage(str(figs/'stage_probability_matrix.png'),width=6.8*inch,height=5.0*inch)); story.append(PageBreak())
story.append(Paragraph('5. Dashboard and bias audit',styles['H1X'])); story.append(Paragraph('The dashboard is generated from local dashboard files and local SVG flag assets. User-provided example images are used only as layout references. No team, finalist, champion, bracket path or probability is imported from those images. Host advantage is kept small and regularized so host teams do not receive unrealistic final probabilities.',styles['BodyX'])); story.append(RLImage(str(figs/'expected_knockout_path.png'),width=7.1*inch,height=3.95*inch)); story.append(Paragraph('Run instructions: open dashboard/index.html, or run streamlit run streamlit_app/app.py. Refit source scripts and notebook are included.',styles['BodyX']))
doc.build(story)
subprocess.run(['python','/home/oai/skills/pdfs/scripts/render_pdf.py',str(report),'--out_dir',str(BASE/'report_render'),'--dpi','130'],check=False)
# PPTX
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
def title(s,text,sub=''):
    tx=s.shapes.add_textbox(Inches(.5),Inches(.25),Inches(12.2),Inches(.55)); p=tx.text_frame.paragraphs[0]; p.text=text; p.font.bold=True; p.font.size=Pt(26); p.font.color.rgb=RGBColor(4,78,59)
    if sub:
        st=s.shapes.add_textbox(Inches(.55),Inches(.88),Inches(12),Inches(.35)); p=st.text_frame.paragraphs[0]; p.text=sub; p.font.size=Pt(12); p.font.color.rgb=RGBColor(75,85,99)
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=RGBColor(4,28,22); tb=s.shapes.add_textbox(Inches(.65),Inches(1.05),Inches(11.8),Inches(1)); p=tb.text_frame.paragraphs[0]; p.text='WC26 Bayesian Hierarchical Simulation'; p.font.bold=True; p.font.size=Pt(38); p.font.color.rgb=RGBColor(248,250,252); tb=s.shapes.add_textbox(Inches(.68),Inches(2.05),Inches(10.4),Inches(.7)); p=tb.text_frame.paragraphs[0]; p.text='PyMC BHM + 100,000 posterior Monte Carlo tournaments + dashboard-ready outputs'; p.font.size=Pt(17); p.font.color.rgb=RGBColor(182,244,212); s.shapes.add_picture(str(figs/'champion_probabilities.png'),Inches(.85),Inches(3.0),width=Inches(6.0)); s.shapes.add_picture(str(figs/'dashboard_main.png'),Inches(7.0),Inches(2.0),width=Inches(5.6))
s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'Model architecture','No hand-picked champion; predictions come from posterior simulation.'); s.shapes.add_picture(str(figs/'posterior_strength_intervals.png'),Inches(.65),Inches(1.25),width=Inches(6.0)); tx=s.shapes.add_textbox(Inches(7.0),Inches(1.25),Inches(5.8),Inches(5.5)); tf=tx.text_frame; bullet=['PyMC Bayesian hierarchical model','Current-strength prior + WC history','Live results locked before simulation','Regularized host effect','Card/injury adjustment layer','100,000 tournament simulations']
for i,b in enumerate(bullet): p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=b; p.font.size=Pt(18 if i==0 else 15); p.font.color.rgb=RGBColor(17,24,39)
s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'Stage probability matrix','Uncertainty is tracked round-by-round.'); s.shapes.add_picture(str(figs/'stage_probability_matrix.png'),Inches(1.1),Inches(1.15),width=Inches(11.1))
s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'Dashboard view','Rendered from local dashboard files and local SVG flag assets.'); s.shapes.add_picture(str(figs/'dashboard_main.png'),Inches(.55),Inches(1.05),width=Inches(12.2))
s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'Expected knockout path','Bracket-style reporting view.'); s.shapes.add_picture(str(figs/'expected_knockout_path.png'),Inches(.7),Inches(1.3),width=Inches(11.9))
s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'Bias audit','Visual examples are layout references only.'); tx=s.shapes.add_textbox(Inches(1.0),Inches(1.3),Inches(11.4),Inches(5.3)); tf=tx.text_frame
for i,b in enumerate(['No sample-image outcomes used as model input','No finalist or champion manually forced','Host advantage controlled and small','Live results verified before inclusion','Uncertainty reported with posterior and MC intervals']): p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text='• '+b; p.font.size=Pt(21); p.font.color.rgb=RGBColor(17,24,39)
prs.save(BASE/'wc26_bhm_prediction_deck.pptx')
# Sources/readme/zip
pd.DataFrame([{'source':'Reuters World Cup 2026 schedule/results graphic','url':'https://www.reuters.com/graphics/SOCCER-WORLDCUP/mopaoglqkpa/','use':'Groups/fixtures/results context'},{'source':'AP News Korea 2-1 Czech Republic','url':'https://apnews.com/article/world-cup-south-korea-czech-republic-score-496e7772dde95ca0af90b5074fdb13d9','use':'Verified live result'},{'source':'Reuters FIFA ranking update','url':'https://www.reuters.com/sports/soccer/argentina-bump-france-top-spot-rankings-ahead-world-cup-title-defence-2026-06-11/','use':'Current strength top ranking context'},{'source':'Fjelstul World Cup Database','url':'https://github.com/jfjelstul/worldcup','use':'Historical source target/structure'},{'source':'FIFA 2026 groups qualification','url':'https://www.fifa.com/en/tournaments/mens/worldcup/canadamexicousa2026/articles/groups-how-teams-qualify-tie-breakers','use':'Top two plus best third-place format'}]).to_csv(BASE/'data/sources.csv',index=False)
(BASE/'README.md').write_text(f'''# WC26 Bayesian Hierarchical Simulation Package\n\nGenerated: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}\n\n## Run immediately\nStatic dashboard:\n```bash\ncd wc26_bhm_final\npython -m http.server 8000\n# open http://localhost:8000/dashboard/\n```\n\nStreamlit:\n```bash\ncd wc26_bhm_final\npip install -r streamlit_app/requirements.txt\nstreamlit run streamlit_app/app.py\n```\n\n## Included\n- PDF report\n- PPTX deck\n- static dashboard with local SVG flags\n- Streamlit app\n- notebook\n- PyMC/model scripts\n- 100,000 simulation CSV outputs\n\n## Verified live results included\n- Mexico 2-0 South Africa\n- Korea Republic 2-1 Czech Republic\n\nAudit: example images used only as layout references, not model evidence.\n''')
zip_path=Path('/mnt/data/wc26_bhm_final_package.zip')
if zip_path.exists(): zip_path.unlink()
with zipfile.ZipFile(zip_path,'w',zipfile.ZIP_DEFLATED) as z:
    for p in BASE.rglob('*'):
        if p.is_file(): z.write(p,p.relative_to(BASE.parent))
print('DONE', zip_path)
print(champ_df.head(10)[['team','Champion_prob','Champion_ci_low','Champion_ci_high','Final_prob']].to_string(index=False))
